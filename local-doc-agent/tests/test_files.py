from __future__ import annotations

import importlib

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

TINY_PNG_BASE64 = (
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+/p9sAAAAASUVORK5CYII="
)


def _reload_modules(monkeypatch, tmp_path):
    monkeypatch.setenv("LOCAL_DOC_AGENT_WORKSPACE", str(tmp_path / "workspace"))
    monkeypatch.setenv("LOCAL_DOC_AGENT_LOGS", str(tmp_path / "logs"))

    import server.config as config
    import server.logging_utils as logging_utils
    import server.tools_assets as tools_assets
    import server.tools_files as tools_files
    import server.tools_templates as tools_templates

    importlib.reload(config)
    importlib.reload(logging_utils)
    importlib.reload(tools_assets)
    importlib.reload(tools_files)
    importlib.reload(tools_templates)
    tools_files.list_templates_tool = tools_templates.list_templates_tool
    tools_files.create_markdown_from_template_tool = tools_templates.create_markdown_from_template_tool
    tools_files.create_docx_from_template_tool = tools_templates.create_docx_from_template_tool
    tools_files.create_pptx_from_template_tool = tools_templates.create_pptx_from_template_tool
    return tools_files, config.settings


def test_write_read_and_patch_text_file(monkeypatch, tmp_path):
    tools, settings = _reload_modules(monkeypatch, tmp_path)

    written = tools.write_text_file_tool("docs/test.md", "# Test\n\nOld text\n")
    assert written["ok"] is True
    assert written["created"] is True

    read = tools.read_text_file_tool("docs/test.md")
    assert read["ok"] is True
    assert "Old text" in read["content"]

    patched = tools.patch_text_file_tool(
        "docs/test.md",
        [{"find": "Old text", "replace": "New text", "must_match_once": True}],
    )
    assert patched["ok"] is True
    assert patched["changed_count"] == 1
    assert patched["backup_path"] is not None

    updated = (settings.workspace_root / "docs" / "test.md").read_text(encoding="utf-8")
    assert "New text" in updated


def test_list_files_filters_extensions(monkeypatch, tmp_path):
    tools, _settings = _reload_modules(monkeypatch, tmp_path)

    tools.write_text_file_tool("docs/a.md", "A")
    tools.write_text_file_tool("docs/b.txt", "B")

    listed = tools.list_files_tool(".", recursive=True, extensions=[".md"])
    assert listed["ok"] is True
    assert listed["files"] == ["docs/a.md"]


def test_create_markdown_from_sections(monkeypatch, tmp_path):
    tools, settings = _reload_modules(monkeypatch, tmp_path)

    created = tools.create_markdown_tool(
        "docs/plan.md",
        title="Project Plan",
        summary="Short summary.",
        sections=[
            {"heading": "Goal", "body": "- Validate local document workflow"},
            {"heading": "Risks", "level": 3, "body": "Connector write permission is unknown."},
        ],
    )

    assert created["ok"] is True
    assert created["created"] is True

    content = (settings.workspace_root / "docs" / "plan.md").read_text(encoding="utf-8")
    assert "# Project Plan" in content
    assert "Short summary." in content
    assert "## Goal" in content
    assert "### Risks" in content


def test_create_markdown_rejects_empty_title(monkeypatch, tmp_path):
    tools, _settings = _reload_modules(monkeypatch, tmp_path)

    created = tools.create_markdown_tool("docs/plan.md", title=" ")

    assert created["ok"] is False
    assert "title" in created["error"]


def test_create_markdown_rejects_suspicious_question_marks(monkeypatch, tmp_path):
    tools, _settings = _reload_modules(monkeypatch, tmp_path)

    created = tools.create_markdown_tool(
        "docs/broken.md",
        title="??? ?? ???",
        summary="깨진 입력",
    )

    assert created["ok"] is False
    assert "인코딩 손상" in created["error"]


def test_export_docx_from_markdown(monkeypatch, tmp_path):
    tools, settings = _reload_modules(monkeypatch, tmp_path)

    tools.write_text_file_tool(
        "docs/source.md",
        "# Proposal\n\nIntro paragraph.\n\n## Scope\n\n- Markdown input\n- DOCX output\n",
    )

    exported = tools.export_docx_from_markdown_tool(
        source_path="docs/source.md",
        output_path="output/source.docx",
    )

    assert exported["ok"] is True
    assert exported["created"] is True
    assert exported["path"] == "output/source.docx"
    assert exported["source_path"] == "docs/source.md"
    assert exported["output_path"] == "output/source.docx"
    assert exported["heading_count"] == 2
    assert exported["paragraph_count"] == 3
    assert exported["style_profile"] == "default_korean_docx_v1"

    document = Document(settings.workspace_root / "output" / "source.docx")
    texts = [paragraph.text for paragraph in document.paragraphs]
    assert "Proposal" in texts
    assert "Intro paragraph." in texts
    assert "Markdown input" in texts
    assert document.styles["Normal"].font.name == "Malgun Gothic"
    assert document.styles["Normal"].font.size.pt == 10.5
    assert round(document.sections[0].left_margin.inches, 2) == 0.85


def test_export_docx_rejects_non_markdown_source(monkeypatch, tmp_path):
    tools, _settings = _reload_modules(monkeypatch, tmp_path)

    result = tools.export_docx_from_markdown_tool(
        source_path="docs/source.txt",
        output_path="output/source.docx",
    )

    assert result["ok"] is False
    assert "입력 파일" in result["error"]


def test_extract_docx_text(monkeypatch, tmp_path):
    tools, settings = _reload_modules(monkeypatch, tmp_path)

    document = Document()
    document.add_heading("전투 시스템 개선안", level=1)
    document.add_paragraph("전투 루프와 성장 구조를 정리한다.")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "항목"
    table.cell(0, 1).text = "상태"
    table.cell(1, 0).text = "DOCX 추출"
    table.cell(1, 1).text = "완료"
    target = settings.workspace_root / "docs" / "source.docx"
    target.parent.mkdir(parents=True, exist_ok=True)
    document.save(target)

    extracted = tools.extract_docx_text_tool(source_path="docs/source.docx")

    assert extracted["ok"] is True
    assert extracted["path"] == "docs/source.docx"
    assert extracted["source_path"] == "docs/source.docx"
    assert extracted["paragraph_count"] == 2
    assert extracted["table_row_count"] == 2
    assert extracted["truncated"] is False
    assert "전투 시스템 개선안" in extracted["text"]
    assert "DOCX 추출 | 완료" in extracted["text"]
    assert extracted["paragraphs"] == ["전투 시스템 개선안", "전투 루프와 성장 구조를 정리한다."]


def test_extract_docx_text_truncates(monkeypatch, tmp_path):
    tools, settings = _reload_modules(monkeypatch, tmp_path)

    document = Document()
    document.add_paragraph("1234567890")
    target = settings.workspace_root / "docs" / "long.docx"
    target.parent.mkdir(parents=True, exist_ok=True)
    document.save(target)

    extracted = tools.extract_docx_text_tool(
        source_path="docs/long.docx",
        max_chars=4,
        include_paragraphs=False,
    )

    assert extracted["ok"] is True
    assert extracted["text"] == "1234"
    assert extracted["truncated"] is True
    assert extracted["paragraphs"] == []


def test_extract_docx_text_rejects_bad_extension(monkeypatch, tmp_path):
    tools, _settings = _reload_modules(monkeypatch, tmp_path)

    result = tools.extract_docx_text_tool(source_path="docs/source.md")

    assert result["ok"] is False
    assert ".docx" in result["error"]


def test_create_xlsx_from_sheets(monkeypatch, tmp_path):
    tools, settings = _reload_modules(monkeypatch, tmp_path)

    created = tools.create_xlsx_from_sheets_tool(
        output_path="output/plan.xlsx",
        sheets=[
            {
                "name": "Schedule",
                "headers": ["Task", "Owner", "Status"],
                "rows": [
                    ["Connector test", "Planner", "Done"],
                    {"Task": "XLSX export", "Owner": "Planner", "Status": "In progress"},
                ],
            },
            {
                "name": "Balance",
                "headers": ["Item", "Value"],
                "rows": [["HP", 100], ["Attack", 12]],
            },
        ],
    )

    assert created["ok"] is True
    assert created["path"] == "output/plan.xlsx"
    assert created["created"] is True
    assert created["sheet_count"] == 2
    assert created["row_count"] == 6

    workbook = load_workbook(settings.workspace_root / "output" / "plan.xlsx")
    assert workbook.sheetnames == ["Schedule", "Balance"]
    assert workbook["Schedule"]["A1"].value == "Task"
    assert workbook["Schedule"]["A3"].value == "XLSX export"
    assert workbook["Balance"]["B3"].value == 12


def test_create_xlsx_creates_backup_on_overwrite(monkeypatch, tmp_path):
    tools, _settings = _reload_modules(monkeypatch, tmp_path)

    first = tools.create_xlsx_from_sheets_tool(
        output_path="output/plan.xlsx",
        sheets=[{"name": "First", "rows": [["A"]]}],
    )
    second = tools.create_xlsx_from_sheets_tool(
        output_path="output/plan.xlsx",
        sheets=[{"name": "Second", "rows": [["B"]]}],
    )

    assert first["ok"] is True
    assert second["ok"] is True
    assert second["created"] is False
    assert second["backup_path"] is not None


def test_create_xlsx_rejects_empty_sheets(monkeypatch, tmp_path):
    tools, _settings = _reload_modules(monkeypatch, tmp_path)

    result = tools.create_xlsx_from_sheets_tool(output_path="output/plan.xlsx", sheets=[])

    assert result["ok"] is False
    assert "sheets" in result["error"]


def test_create_xlsx_rejects_suspicious_question_marks(monkeypatch, tmp_path):
    tools, _settings = _reload_modules(monkeypatch, tmp_path)

    result = tools.create_xlsx_from_sheets_tool(
        output_path="output/broken.xlsx",
        sheets=[
            {
                "name": "Checklist",
                "rows": [["MCP ??", "Planner", "Done"]],
            }
        ],
    )

    assert result["ok"] is False
    assert "인코딩 손상" in result["error"]


def test_extract_xlsx_text(monkeypatch, tmp_path):
    tools, _settings = _reload_modules(monkeypatch, tmp_path)

    tools.create_xlsx_from_sheets_tool(
        output_path="output/plan.xlsx",
        sheets=[
            {
                "name": "Checklist",
                "headers": ["항목", "담당", "상태"],
                "rows": [
                    ["DOCX 검증", "Planner", "완료"],
                    ["XLSX 검색", "Planner", "진행 중"],
                ],
            }
        ],
    )

    extracted = tools.extract_xlsx_text_tool(
        source_path="output/plan.xlsx",
        query="검증",
    )

    assert extracted["ok"] is True
    assert extracted["path"] == "output/plan.xlsx"
    assert extracted["sheet_count"] == 1
    assert extracted["cell_count"] == 9
    assert extracted["match_count"] == 1
    assert extracted["matches"][0]["sheet"] == "Checklist"
    assert extracted["matches"][0]["coordinate"] == "A2"
    assert "Checklist!A2: DOCX 검증" in extracted["text"]


def test_extract_xlsx_text_truncates_cells(monkeypatch, tmp_path):
    tools, _settings = _reload_modules(monkeypatch, tmp_path)

    tools.create_xlsx_from_sheets_tool(
        output_path="output/plan.xlsx",
        sheets=[{"name": "Data", "rows": [["A"], ["B"], ["C"]]}],
    )

    extracted = tools.extract_xlsx_text_tool(
        source_path="output/plan.xlsx",
        max_cells=2,
    )

    assert extracted["ok"] is True
    assert extracted["cell_count"] == 2
    assert extracted["truncated"] is True


def test_extract_xlsx_text_rejects_bad_extension(monkeypatch, tmp_path):
    tools, _settings = _reload_modules(monkeypatch, tmp_path)

    result = tools.extract_xlsx_text_tool(source_path="output/plan.csv")

    assert result["ok"] is False
    assert ".xlsx" in result["error"]


def test_create_pptx_from_spec(monkeypatch, tmp_path):
    tools, settings = _reload_modules(monkeypatch, tmp_path)

    created = tools.create_pptx_from_spec_tool(
        output_path="output/deck.pptx",
        title="MVP Review",
        subtitle="Local document agent",
        slides=[
            {
                "title": "Goals",
                "bullets": ["Validate MCP calls", "Generate editable files"],
                "notes": "Use as a draft deck.",
            },
            {
                "title": "Next Steps",
                "body": "Add templates after conversion tools stabilize.",
            },
        ],
    )

    assert created["ok"] is True
    assert created["path"] == "output/deck.pptx"
    assert created["created"] is True
    assert created["slide_count"] == 3
    assert created["bullet_count"] == 2
    assert created["note_count"] == 1
    assert created["style_profile"] == "default_korean_pptx_v1"

    presentation = Presentation(settings.workspace_root / "output" / "deck.pptx")
    assert len(presentation.slides) == 3
    assert round(presentation.slide_width.inches, 2) == 13.33
    assert round(presentation.slide_height.inches, 2) == 7.5
    assert presentation.slides[0].shapes.title.text == "MVP Review"
    assert presentation.slides[0].shapes.title.text_frame.paragraphs[0].font.name == "Malgun Gothic"
    slide_text = "\n".join(shape.text for shape in presentation.slides[1].shapes if hasattr(shape, "text"))
    assert "Validate MCP calls" in slide_text


def test_create_pptx_creates_backup_on_overwrite(monkeypatch, tmp_path):
    tools, _settings = _reload_modules(monkeypatch, tmp_path)

    first = tools.create_pptx_from_spec_tool(
        output_path="output/deck.pptx",
        title="First",
        slides=[{"title": "One", "bullets": ["A"]}],
    )
    second = tools.create_pptx_from_spec_tool(
        output_path="output/deck.pptx",
        title="Second",
        slides=[{"title": "Two", "bullets": ["B"]}],
    )

    assert first["ok"] is True
    assert second["ok"] is True
    assert second["created"] is False
    assert second["backup_path"] is not None


def test_create_pptx_rejects_empty_slides(monkeypatch, tmp_path):
    tools, _settings = _reload_modules(monkeypatch, tmp_path)

    result = tools.create_pptx_from_spec_tool(output_path="output/deck.pptx", title="Deck", slides=[])

    assert result["ok"] is False
    assert "slides" in result["error"]


def test_create_pptx_rejects_suspicious_question_marks(monkeypatch, tmp_path):
    tools, _settings = _reload_modules(monkeypatch, tmp_path)

    result = tools.create_pptx_from_spec_tool(
        output_path="output/broken.pptx",
        title="??? DOCX ???",
        slides=[{"title": "Slide", "body": "Body"}],
    )

    assert result["ok"] is False
    assert "인코딩 손상" in result["error"]


def test_save_base64_image_and_list_assets(monkeypatch, tmp_path):
    tools, settings = _reload_modules(monkeypatch, tmp_path)

    saved = tools.save_base64_image_tool(
        output_path="assets/pixel.png",
        image_base64=TINY_PNG_BASE64,
    )

    assert saved["ok"] is True
    assert saved["path"] == "assets/pixel.png"
    assert saved["created"] is True
    assert saved["byte_count"] > 0

    listed = tools.list_assets_tool()
    assert listed["ok"] is True
    assert listed["files"] == ["assets/pixel.png"]
    assert (settings.workspace_root / "assets" / "pixel.png").exists()


def test_read_image_file(monkeypatch, tmp_path):
    tools, _settings = _reload_modules(monkeypatch, tmp_path)

    tools.save_base64_image_tool(
        output_path="assets/pixel.png",
        image_base64=TINY_PNG_BASE64,
    )

    read = tools.read_image_file_tool(
        path="assets/pixel.png",
        include_base64=True,
        include_data_uri=True,
    )

    assert read["ok"] is True
    assert read["path"] == "assets/pixel.png"
    assert read["mime_type"] == "image/png"
    assert read["width"] == 1
    assert read["height"] == 1
    assert read["byte_count"] > 0
    assert read["image_base64"] == TINY_PNG_BASE64
    assert read["data_uri"].startswith("data:image/png;base64,")


def test_read_image_file_respects_max_bytes(monkeypatch, tmp_path):
    tools, _settings = _reload_modules(monkeypatch, tmp_path)

    tools.save_base64_image_tool(
        output_path="assets/pixel.png",
        image_base64=TINY_PNG_BASE64,
    )

    read = tools.read_image_file_tool(
        path="assets/pixel.png",
        max_bytes=1,
    )

    assert read["ok"] is False
    assert "max_bytes" in read["error"]


def test_read_image_file_rejects_bad_extension(monkeypatch, tmp_path):
    tools, _settings = _reload_modules(monkeypatch, tmp_path)

    read = tools.read_image_file_tool(path="docs/source.md")

    assert read["ok"] is False
    assert ".png" in read["error"]


def test_insert_image_to_markdown(monkeypatch, tmp_path):
    tools, settings = _reload_modules(monkeypatch, tmp_path)

    tools.write_text_file_tool("docs/visual.md", "# Visual\n")
    tools.save_base64_image_tool("assets/pixel.png", TINY_PNG_BASE64)

    inserted = tools.insert_image_to_markdown_tool(
        markdown_path="docs/visual.md",
        image_path="assets/pixel.png",
        alt_text="Pixel",
    )

    assert inserted["ok"] is True
    assert inserted["path"] == "docs/visual.md"
    assert inserted["image_path"] == "assets/pixel.png"
    assert inserted["backup_path"] is not None

    content = (settings.workspace_root / "docs" / "visual.md").read_text(encoding="utf-8")
    assert "![Pixel](../assets/pixel.png)" in content


def test_save_base64_image_rejects_invalid_base64(monkeypatch, tmp_path):
    tools, _settings = _reload_modules(monkeypatch, tmp_path)

    result = tools.save_base64_image_tool("assets/pixel.png", "not-base64")

    assert result["ok"] is False
    assert "base64" in result["error"]


def test_insert_image_to_pptx(monkeypatch, tmp_path):
    tools, settings = _reload_modules(monkeypatch, tmp_path)

    tools.create_pptx_from_spec_tool(
        output_path="output/deck.pptx",
        title="Image Deck",
        slides=[{"title": "Image Slide", "bullets": ["Before image"]}],
    )
    tools.save_base64_image_tool("assets/pixel.png", TINY_PNG_BASE64)

    inserted = tools.insert_image_to_pptx_tool(
        pptx_path="output/deck.pptx",
        image_path="assets/pixel.png",
        slide_index=1,
        left=1.0,
        top=2.0,
        width=1.0,
    )

    assert inserted["ok"] is True
    assert inserted["path"] == "output/deck.pptx"
    assert inserted["image_path"] == "assets/pixel.png"
    assert inserted["backup_path"] is not None
    assert inserted["slide_index"] == 1

    presentation = Presentation(settings.workspace_root / "output" / "deck.pptx")
    picture_shapes = [shape for shape in presentation.slides[1].shapes if shape.shape_type == 13]
    assert len(picture_shapes) == 1


def test_insert_image_to_pptx_rejects_bad_slide_index(monkeypatch, tmp_path):
    tools, _settings = _reload_modules(monkeypatch, tmp_path)

    tools.create_pptx_from_spec_tool(
        output_path="output/deck.pptx",
        title="Image Deck",
        slides=[{"title": "Image Slide", "bullets": ["Before image"]}],
    )
    tools.save_base64_image_tool("assets/pixel.png", TINY_PNG_BASE64)

    result = tools.insert_image_to_pptx_tool(
        pptx_path="output/deck.pptx",
        image_path="assets/pixel.png",
        slide_index=99,
    )

    assert result["ok"] is False
    assert "slide_index" in result["error"]


def test_list_templates(monkeypatch, tmp_path):
    tools, _settings = _reload_modules(monkeypatch, tmp_path)

    result = tools.list_templates_tool()

    assert result["ok"] is True
    names = {template["name"] for template in result["templates"]}
    assert {"planning_doc", "proposal_doc", "checklist_doc"}.issubset(names)


def test_create_markdown_from_template(monkeypatch, tmp_path):
    tools, settings = _reload_modules(monkeypatch, tmp_path)

    created = tools.create_markdown_from_template_tool(
        template_name="planning_doc",
        output_path="docs/template-plan.md",
        title="전투 시스템 기획서",
        summary="전투 루프와 성장 구조를 정리한다.",
    )

    assert created["ok"] is True
    assert created["path"] == "docs/template-plan.md"
    assert created["template_name"] == "planning_doc"
    assert created["section_count"] > 0

    content = (settings.workspace_root / "docs" / "template-plan.md").read_text(encoding="utf-8")
    assert "# 전투 시스템 기획서" in content
    assert "## 핵심 루프" in content
    assert "전투 루프와 성장 구조를 정리한다." in content


def test_create_markdown_from_template_creates_backup(monkeypatch, tmp_path):
    tools, _settings = _reload_modules(monkeypatch, tmp_path)

    first = tools.create_markdown_from_template_tool(
        template_name="checklist_doc",
        output_path="docs/checklist.md",
        title="첫 체크리스트",
    )
    second = tools.create_markdown_from_template_tool(
        template_name="checklist_doc",
        output_path="docs/checklist.md",
        title="두 번째 체크리스트",
    )

    assert first["ok"] is True
    assert second["ok"] is True
    assert second["created"] is False
    assert second["backup_path"] is not None


def test_create_markdown_from_template_rejects_unknown_template(monkeypatch, tmp_path):
    tools, _settings = _reload_modules(monkeypatch, tmp_path)

    result = tools.create_markdown_from_template_tool(
        template_name="unknown",
        output_path="docs/unknown.md",
        title="Unknown",
    )

    assert result["ok"] is False
    assert "template_name" in result["error"]


def test_create_docx_from_template(monkeypatch, tmp_path):
    tools, settings = _reload_modules(monkeypatch, tmp_path)

    created = tools.create_docx_from_template_tool(
        template_name="proposal_doc",
        output_path="output/proposal.docx",
        title="문서 자동화 제안서",
        summary="로컬 문서 작업 자동화를 제안한다.",
    )

    assert created["ok"] is True
    assert created["path"] == "output/proposal.docx"
    assert created["output_path"] == "output/proposal.docx"
    assert created["template_name"] == "proposal_doc"
    assert created["heading_count"] > 0
    assert created["paragraph_count"] > 0
    assert created["style_profile"] == "default_korean_docx_v1"

    document = Document(settings.workspace_root / "output" / "proposal.docx")
    texts = [paragraph.text for paragraph in document.paragraphs]
    assert "문서 자동화 제안서" in texts
    assert "제안 배경" in texts
    assert "로컬 문서 작업 자동화를 제안한다." in texts
    assert document.styles["Heading 1"].font.name == "Malgun Gothic"


def test_create_docx_from_template_creates_backup(monkeypatch, tmp_path):
    tools, _settings = _reload_modules(monkeypatch, tmp_path)

    first = tools.create_docx_from_template_tool(
        template_name="planning_doc",
        output_path="output/template.docx",
        title="첫 문서",
    )
    second = tools.create_docx_from_template_tool(
        template_name="planning_doc",
        output_path="output/template.docx",
        title="두 번째 문서",
    )

    assert first["ok"] is True
    assert second["ok"] is True
    assert second["created"] is False
    assert second["backup_path"] is not None


def test_create_docx_from_template_rejects_bad_extension(monkeypatch, tmp_path):
    tools, _settings = _reload_modules(monkeypatch, tmp_path)

    result = tools.create_docx_from_template_tool(
        template_name="planning_doc",
        output_path="output/template.md",
        title="잘못된 확장자",
    )

    assert result["ok"] is False
    assert ".docx" in result["error"]


def test_create_docx_from_template_rejects_suspicious_question_marks(monkeypatch, tmp_path):
    tools, _settings = _reload_modules(monkeypatch, tmp_path)

    result = tools.create_docx_from_template_tool(
        template_name="proposal_doc",
        output_path="output/broken.docx",
        title="??? DOCX ???",
    )

    assert result["ok"] is False
    assert "인코딩 손상" in result["error"]


def test_create_pptx_from_template(monkeypatch, tmp_path):
    tools, settings = _reload_modules(monkeypatch, tmp_path)

    created = tools.create_pptx_from_template_tool(
        template_name="planning_doc",
        output_path="output/planning.pptx",
        title="전투 시스템 개선안",
        summary="전투 루프와 성장 구조를 정리한다.",
    )

    assert created["ok"] is True
    assert created["path"] == "output/planning.pptx"
    assert created["output_path"] == "output/planning.pptx"
    assert created["template_name"] == "planning_doc"
    assert created["slide_count"] == 7
    assert created["bullet_count"] > 0
    assert created["body_count"] > 0
    assert created["style_profile"] == "default_korean_pptx_v1"

    presentation = Presentation(settings.workspace_root / "output" / "planning.pptx")
    assert round(presentation.slide_width.inches, 2) == 13.33
    slide_texts = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_texts.append(shape.text)

    assert any("전투 시스템 개선안" in text for text in slide_texts)
    assert any("개요" in text for text in slide_texts)
    assert any("핵심 루프" in text for text in slide_texts)


def test_create_pptx_from_template_creates_backup(monkeypatch, tmp_path):
    tools, _settings = _reload_modules(monkeypatch, tmp_path)

    first = tools.create_pptx_from_template_tool(
        template_name="proposal_doc",
        output_path="output/template.pptx",
        title="첫 발표자료",
    )
    second = tools.create_pptx_from_template_tool(
        template_name="proposal_doc",
        output_path="output/template.pptx",
        title="두 번째 발표자료",
    )

    assert first["ok"] is True
    assert second["ok"] is True
    assert second["created"] is False
    assert second["backup_path"] is not None


def test_create_pptx_from_template_rejects_bad_extension(monkeypatch, tmp_path):
    tools, _settings = _reload_modules(monkeypatch, tmp_path)

    result = tools.create_pptx_from_template_tool(
        template_name="planning_doc",
        output_path="output/template.docx",
        title="잘못된 확장자",
    )

    assert result["ok"] is False
    assert ".pptx" in result["error"]
