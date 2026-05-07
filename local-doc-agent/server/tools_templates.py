from __future__ import annotations

import difflib
import re
import shutil
from datetime import datetime
from pathlib import Path
from typing import Any

from .config import ensure_base_directories, settings
from .docx_style import apply_docx_style
from .logging_utils import write_operation_log
from .pptx_style import apply_pptx_style, style_content_slide, style_paragraph, style_title_slide


TEMPLATES: dict[str, dict[str, Any]] = {
    "planning_doc": {
        "name": "planning_doc",
        "title": "{{title}}",
        "description": "게임 기획서 초안",
        "sections": [
            ("개요", "{{summary}}"),
            ("목표", "- 핵심 목표\n- 사용자 가치\n- 완료 기준"),
            ("핵심 루프", "1. 진입\n2. 선택\n3. 결과\n4. 보상\n5. 재시도"),
            ("시스템 구성", "- 전투\n- 성장\n- 콘텐츠\n- UX\n- 밸런스"),
            ("리스크", "- UX 리스크\n- 밸런스 리스크\n- 개발 비용"),
            ("다음 작업", "- 세부 규칙 정의\n- 예외 처리 정리\n- 테스트 시나리오 작성"),
        ],
    },
    "proposal_doc": {
        "name": "proposal_doc",
        "title": "{{title}}",
        "description": "제안서 초안",
        "sections": [
            ("제안 배경", "{{summary}}"),
            ("문제 정의", "- 현재 문제\n- 영향 범위\n- 해결 필요성"),
            ("제안 내용", "- 핵심 제안\n- 적용 방식\n- 기대 효과"),
            ("범위", "- 포함 범위\n- 제외 범위\n- 후속 검토"),
            ("일정", "- 1단계\n- 2단계\n- 검수"),
            ("결론", "실행 가능성과 기대 효과를 기준으로 우선순위를 판단한다."),
        ],
    },
    "checklist_doc": {
        "name": "checklist_doc",
        "title": "{{title}}",
        "description": "검증 체크리스트 초안",
        "sections": [
            ("검증 대상", "{{summary}}"),
            ("사전 조건", "- [ ] 환경 준비\n- [ ] 파일 경로 확인\n- [ ] 백업 정책 확인"),
            ("검증 항목", "- [ ] 정상 케이스\n- [ ] 오류 케이스\n- [ ] 경계값\n- [ ] 로그 기록"),
            ("결과 기록", "- 완료 항목\n- 실패 항목\n- 재검증 필요 항목"),
            ("후속 작업", "- 수정 필요 사항\n- 추가 테스트\n- 문서 보강"),
        ],
    },
}


def _to_workspace_relative(path: Path) -> str:
    return path.resolve().relative_to(settings.workspace_root).as_posix()


def _resolve_workspace_path(path: str | None) -> Path:
    raw = "." if path is None or path == "" else path
    candidate = (settings.workspace_root / raw).resolve()
    try:
        candidate.relative_to(settings.workspace_root)
    except ValueError as exc:
        raise ValueError(f"workspace 밖의 경로는 사용할 수 없음: {path}") from exc
    return candidate


def _diff_summary(before: str, after: str) -> dict[str, int]:
    before_lines = before.splitlines(keepends=True)
    after_lines = after.splitlines(keepends=True)
    diff = difflib.unified_diff(before_lines, after_lines, lineterm="")
    added = 0
    removed = 0
    for line in diff:
        if line.startswith("+++") or line.startswith("---"):
            continue
        if line.startswith("+"):
            added += 1
        elif line.startswith("-"):
            removed += 1
    return {"added_lines": added, "removed_lines": removed}


def _backup_existing_file(path: Path) -> str | None:
    if not path.exists() or not path.is_file():
        return None

    relative = path.relative_to(settings.workspace_root)
    backup_dir = settings.backups_root / relative.parent
    backup_dir.mkdir(parents=True, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    backup_name = f"{path.stem}.{timestamp}{path.suffix}"
    backup_path = backup_dir / backup_name

    counter = 1
    while backup_path.exists():
        backup_path = backup_dir / f"{path.stem}.{timestamp}-{counter}{path.suffix}"
        counter += 1

    shutil.copy2(path, backup_path)
    return _to_workspace_relative(backup_path)


def _safe_result(tool_name: str, action: Any) -> dict[str, Any]:
    try:
        ensure_base_directories()
        result = action()
    except Exception as exc:  # Tool boundaries should return readable errors to ChatGPT.
        result = {"ok": False, "error": str(exc)}
    write_operation_log(tool_name, result)
    return result


def _render_template_text(value: str, variables: dict[str, Any]) -> str:
    def replace(match: re.Match[str]) -> str:
        key = match.group(1).strip()
        return str(variables.get(key, ""))

    return re.sub(r"\{\{\s*([a-zA-Z0-9_]+)\s*\}\}", replace, value)


def _build_template_markdown(
    template: dict[str, Any],
    title: str,
    summary: str,
    variables: dict[str, Any],
) -> str:
    merged_variables = {"title": title, "summary": summary, **variables}
    rendered_title = _render_template_text(str(template["title"]), merged_variables).strip()
    if not rendered_title:
        raise ValueError("title 값은 비어 있을 수 없음")

    lines = [f"# {rendered_title}", ""]
    for heading, body in template["sections"]:
        rendered_heading = _render_template_text(heading, merged_variables).strip()
        rendered_body = _render_template_text(body, merged_variables).strip()
        if rendered_heading:
            lines.extend([f"## {rendered_heading}", ""])
        if rendered_body:
            lines.extend([rendered_body, ""])
    return "\n".join(lines).rstrip() + "\n"


def _add_markdown_to_docx(document: Any, markdown: str) -> dict[str, int]:
    heading_count = 0
    paragraph_count = 0

    for raw_line in markdown.splitlines():
        line = raw_line.rstrip()
        stripped = line.strip()
        if not stripped:
            continue

        if stripped.startswith("#"):
            marker, _, heading = stripped.partition(" ")
            if marker and set(marker) == {"#"} and heading.strip():
                level = min(max(len(marker), 1), 4)
                document.add_heading(heading.strip(), level=level)
                heading_count += 1
                continue

        if stripped.startswith("- [ ] "):
            document.add_paragraph(stripped[6:].strip(), style="List Bullet")
            paragraph_count += 1
            continue

        if stripped.startswith("- "):
            document.add_paragraph(stripped[2:].strip(), style="List Bullet")
            paragraph_count += 1
            continue

        if re.match(r"^\d+\.\s+", stripped):
            document.add_paragraph(re.sub(r"^\d+\.\s+", "", stripped), style="List Number")
            paragraph_count += 1
            continue

        document.add_paragraph(stripped)
        paragraph_count += 1

    return {"heading_count": heading_count, "paragraph_count": paragraph_count}


def _section_body_to_slide_parts(body: str) -> tuple[str, list[str]]:
    lines = [line.strip() for line in body.splitlines() if line.strip()]
    bullets: list[str] = []
    body_lines: list[str] = []

    for line in lines:
        if line.startswith("- [ ] "):
            bullets.append(line[6:].strip())
        elif line.startswith("- "):
            bullets.append(line[2:].strip())
        elif re.match(r"^\d+\.\s+", line):
            bullets.append(re.sub(r"^\d+\.\s+", "", line).strip())
        else:
            body_lines.append(line)

    return "\n".join(body_lines), bullets


def _build_template_slides(
    template: dict[str, Any],
    title: str,
    summary: str,
    variables: dict[str, Any],
) -> tuple[str, list[dict[str, Any]]]:
    merged_variables = {"title": title, "summary": summary, **variables}
    rendered_title = _render_template_text(str(template["title"]), merged_variables).strip()
    if not rendered_title:
        raise ValueError("title 값은 비어 있을 수 없음")

    slides: list[dict[str, Any]] = []
    for heading, body in template["sections"]:
        rendered_heading = _render_template_text(heading, merged_variables).strip()
        rendered_body = _render_template_text(body, merged_variables).strip()
        if not rendered_heading and not rendered_body:
            continue

        body_text, bullets = _section_body_to_slide_parts(rendered_body)
        slides.append(
            {
                "title": rendered_heading or rendered_title,
                "body": body_text,
                "bullets": bullets,
            }
        )

    return rendered_title, slides


def _write_template_pptx_result(
    template: dict[str, Any],
    template_name: str,
    output_path: str,
    title: str,
    summary: str,
    variables: dict[str, Any],
    overwrite: bool | None = None,
    create_backup: bool | None = None,
) -> dict[str, Any]:
    from pptx import Presentation
    rendered_title, slides = _build_template_slides(
        template=template,
        title=title,
        summary=summary,
        variables=variables,
    )
    if not slides:
        raise ValueError("템플릿에 생성 가능한 슬라이드가 없음")

    target = _resolve_workspace_path(output_path)
    if target.suffix.lower() != ".pptx":
        raise ValueError("템플릿 PPTX 출력 파일은 .pptx 확장자만 지원함")

    should_overwrite = settings.default_overwrite if overwrite is None else overwrite
    should_backup = settings.create_backup if create_backup is None else create_backup

    exists = target.exists()
    if exists and not should_overwrite:
        raise FileExistsError(f"파일이 이미 존재함: {output_path}")
    if exists and not target.is_file():
        raise ValueError(f"파일이 아님: {output_path}")

    presentation = Presentation()
    style_profile = apply_pptx_style(presentation)
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = rendered_title
    try:
        title_slide.placeholders[1].text = summary.strip()
    except IndexError:
        pass
    style_title_slide(title_slide)

    bullet_count = 0
    body_count = 0

    for slide_spec in slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        if slide.shapes.title:
            slide.shapes.title.text = slide_spec["title"]
        style_content_slide(slide)

        content = slide.placeholders[1].text_frame
        content.clear()

        body = slide_spec["body"]
        bullets = slide_spec["bullets"]
        if body:
            paragraph = content.paragraphs[0]
            paragraph.text = body
            style_paragraph(paragraph, size_pt=20)
            body_count += 1

        for index, bullet in enumerate(bullets):
            paragraph = content.paragraphs[0] if index == 0 and not body else content.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            style_paragraph(paragraph, size_pt=20)
            bullet_count += 1

    backup_path = _backup_existing_file(target) if exists and should_backup else None
    target.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(target)

    return {
        "ok": True,
        "path": _to_workspace_relative(target),
        "output_path": _to_workspace_relative(target),
        "created": not exists,
        "backup_path": backup_path,
        "template_name": template_name,
        "section_count": len(template["sections"]),
        "slide_count": len(presentation.slides),
        "bullet_count": bullet_count,
        "body_count": body_count,
        "style_profile": style_profile,
    }


def list_templates_tool() -> dict[str, Any]:
    def action() -> dict[str, Any]:
        templates = [
            {
                "name": template["name"],
                "description": template["description"],
                "section_count": len(template["sections"]),
            }
            for template in TEMPLATES.values()
        ]
        return {"ok": True, "templates": templates}

    return _safe_result("list_templates", action)


def create_markdown_from_template_tool(
    template_name: str,
    output_path: str,
    title: str,
    summary: str = "",
    variables: dict[str, Any] | None = None,
    overwrite: bool | None = None,
    create_backup: bool | None = None,
) -> dict[str, Any]:
    def action() -> dict[str, Any]:
        template = TEMPLATES.get(template_name)
        if template is None:
            allowed = ", ".join(sorted(TEMPLATES))
            raise ValueError(f"알 수 없는 template_name: {template_name}. 사용 가능: {allowed}")

        target = _resolve_workspace_path(output_path)
        if target.suffix.lower() != ".md":
            raise ValueError("템플릿 출력 파일은 .md 확장자만 지원함")

        should_overwrite = settings.default_overwrite if overwrite is None else overwrite
        should_backup = settings.create_backup if create_backup is None else create_backup

        exists = target.exists()
        if exists and not should_overwrite:
            raise FileExistsError(f"파일이 이미 존재함: {output_path}")
        if exists and not target.is_file():
            raise ValueError(f"파일이 아님: {output_path}")

        content = _build_template_markdown(
            template=template,
            title=title,
            summary=summary,
            variables=variables or {},
        )
        before = target.read_text(encoding="utf-8") if exists else ""
        backup_path = _backup_existing_file(target) if exists and should_backup else None

        target.parent.mkdir(parents=True, exist_ok=True)
        target.write_text(content, encoding="utf-8", newline="\n")

        return {
            "ok": True,
            "path": _to_workspace_relative(target),
            "created": not exists,
            "backup_path": backup_path,
            "template_name": template_name,
            "section_count": len(template["sections"]),
            "diff_summary": _diff_summary(before, content),
        }

    return _safe_result("create_markdown_from_template", action)


def create_docx_from_template_tool(
    template_name: str,
    output_path: str,
    title: str,
    summary: str = "",
    variables: dict[str, Any] | None = None,
    overwrite: bool | None = None,
    create_backup: bool | None = None,
) -> dict[str, Any]:
    def action() -> dict[str, Any]:
        from docx import Document

        template = TEMPLATES.get(template_name)
        if template is None:
            allowed = ", ".join(sorted(TEMPLATES))
            raise ValueError(f"알 수 없는 template_name: {template_name}. 사용 가능: {allowed}")

        target = _resolve_workspace_path(output_path)
        if target.suffix.lower() != ".docx":
            raise ValueError("템플릿 DOCX 출력 파일은 .docx 확장자만 지원함")

        should_overwrite = settings.default_overwrite if overwrite is None else overwrite
        should_backup = settings.create_backup if create_backup is None else create_backup

        exists = target.exists()
        if exists and not should_overwrite:
            raise FileExistsError(f"파일이 이미 존재함: {output_path}")
        if exists and not target.is_file():
            raise ValueError(f"파일이 아님: {output_path}")

        markdown = _build_template_markdown(
            template=template,
            title=title,
            summary=summary,
            variables=variables or {},
        )
        document = Document()
        style_profile = apply_docx_style(document)
        counts = _add_markdown_to_docx(document, markdown)

        backup_path = _backup_existing_file(target) if exists and should_backup else None
        target.parent.mkdir(parents=True, exist_ok=True)
        document.save(target)

        return {
            "ok": True,
            "path": _to_workspace_relative(target),
            "output_path": _to_workspace_relative(target),
            "created": not exists,
            "backup_path": backup_path,
            "template_name": template_name,
            "section_count": len(template["sections"]),
            "style_profile": style_profile,
            **counts,
        }

    return _safe_result("create_docx_from_template", action)


def create_pptx_from_template_tool(
    template_name: str,
    output_path: str,
    title: str,
    summary: str = "",
    variables: dict[str, Any] | None = None,
    overwrite: bool | None = None,
    create_backup: bool | None = None,
) -> dict[str, Any]:
    def action() -> dict[str, Any]:
        template = TEMPLATES.get(template_name)
        if template is None:
            allowed = ", ".join(sorted(TEMPLATES))
            raise ValueError(f"알 수 없는 template_name: {template_name}. 사용 가능: {allowed}")

        return _write_template_pptx_result(
            template=template,
            template_name=template_name,
            output_path=output_path,
            title=title,
            summary=summary,
            variables=variables or {},
            overwrite=overwrite,
            create_backup=create_backup,
        )

    return _safe_result("create_pptx_from_template", action)
