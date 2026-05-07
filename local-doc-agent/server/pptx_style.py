from __future__ import annotations

from typing import Any


PPTX_STYLE_PROFILE = "default_korean_pptx_v1"
PPTX_BODY_FONT = "Malgun Gothic"


def apply_pptx_style(presentation: Any) -> str:
    from pptx.util import Inches

    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    return PPTX_STYLE_PROFILE


def style_paragraph(paragraph: Any, *, size_pt: int, bold: bool = False) -> None:
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    paragraph.font.name = PPTX_BODY_FONT
    paragraph.font.size = Pt(size_pt)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = RGBColor(36, 36, 36)


def style_text_frame(text_frame: Any, *, size_pt: int, bold: bool = False) -> None:
    text_frame.word_wrap = True
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    for paragraph in text_frame.paragraphs:
        style_paragraph(paragraph, size_pt=size_pt, bold=bold)


def style_title_slide(slide: Any) -> None:
    if slide.shapes.title:
        style_text_frame(slide.shapes.title.text_frame, size_pt=34, bold=True)

    try:
        subtitle = slide.placeholders[1]
    except IndexError:
        return
    if hasattr(subtitle, "text_frame"):
        style_text_frame(subtitle.text_frame, size_pt=18)


def style_content_slide(slide: Any) -> None:
    if slide.shapes.title:
        style_text_frame(slide.shapes.title.text_frame, size_pt=26, bold=True)
