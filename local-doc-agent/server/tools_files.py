from __future__ import annotations

import difflib
import shutil
from datetime import datetime
from pathlib import Path
from typing import Any

from .config import ensure_base_directories, settings
from .docx_style import apply_docx_style
from .logging_utils import write_operation_log
from .pptx_style import apply_pptx_style, style_content_slide, style_paragraph, style_title_slide
from .text_guard import reject_suspicious_question_marks
from .tools_assets import (
    insert_image_to_markdown_tool,
    insert_image_to_pptx_tool,
    list_assets_tool,
    save_base64_image_tool,
)


TEXT_EXTENSIONS = {".md", ".txt", ".json", ".csv", ".yaml", ".yml"}
DEFAULT_LIST_EXTENSIONS = [".md", ".txt", ".docx", ".pptx", ".xlsx", ".png", ".jpg", ".jpeg", ".webp"]


def _to_workspace_relative(path: Path) -> str:
    return path.resolve().relative_to(settings.workspace_root).as_posix()


def _resolve_workspace_path(path: str | None) -> Path:
    raw = "." if path is None or path == "" else path
    candidate = (settings.workspace_root / raw).resolve()
    try:
        candidate.relative_to(settings.workspace_root)
    except ValueError as exc:
        raise ValueError(f"workspace 밖의 경로는 사용할 수 없음: {path}") from exc
    return candidate


def _ensure_text_extension(path: Path) -> None:
    if path.suffix.lower() not in TEXT_EXTENSIONS:
        allowed = ", ".join(sorted(TEXT_EXTENSIONS))
        raise ValueError(f"텍스트 도구는 다음 확장자만 지원함: {allowed}")


def _ensure_extension(path: Path, allowed_extensions: set[str], label: str) -> None:
    if path.suffix.lower() not in allowed_extensions:
        allowed = ", ".join(sorted(allowed_extensions))
        raise ValueError(f"{label}은 다음 확장자만 지원함: {allowed}")


def _diff_summary(before: str, after: str) -> dict[str, int]:
    before_lines = before.splitlines(keepends=True)
    after_lines = after.splitlines(keepends=True)
    diff = difflib.unified_diff(before_lines, after_lines, lineterm="")
    added = 0
    removed = 0
    for line in diff:
        if line.startswith("+++") or line.startswith("---"):
            continue
        if line.startswith("+"):
            added += 1
        elif line.startswith("-"):
            removed += 1
    return {"added_lines": added, "removed_lines": removed}


def _normalize_markdown_level(value: Any) -> int:
    try:
        level = int(value)
    except (TypeError, ValueError):
        level = 2
    return min(max(level, 2), 6)


def _build_markdown_content(
    title: str,
    summary: str | None = None,
    sections: list[dict[str, Any]] | None = None,
) -> str:
    reject_suspicious_question_marks(title, label="title")
    reject_suspicious_question_marks(summary, label="summary")
    reject_suspicious_question_marks(sections, label="sections")

    normalized_title = title.strip()
    if not normalized_title:
        raise ValueError("title 값은 비어 있을 수 없음")

    lines = [f"# {normalized_title}", ""]

    if summary is not None and summary.strip():
        lines.extend([summary.strip(), ""])

    for section in sections or []:
        heading = str(section.get("heading", "")).strip()
        body = str(section.get("body", "")).strip()
        if not heading and not body:
            continue
        if heading:
            level = _normalize_markdown_level(section.get("level", 2))
            lines.extend([f"{'#' * level} {heading}", ""])
        if body:
            lines.extend([body, ""])

    return "\n".join(lines).rstrip() + "\n"


def _backup_existing_file(path: Path) -> str | None:
    if not path.exists() or not path.is_file():
        return None

    relative = path.relative_to(settings.workspace_root)
    backup_dir = settings.backups_root / relative.parent
    backup_dir.mkdir(parents=True, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    backup_name = f"{path.stem}.{timestamp}{path.suffix}"
    backup_path = backup_dir / backup_name

    counter = 1
    while backup_path.exists():
        backup_path = backup_dir / f"{path.stem}.{timestamp}-{counter}{path.suffix}"
        counter += 1

    shutil.copy2(path, backup_path)
    return _to_workspace_relative(backup_path)


def _write_text_file_result(
    path: str,
    content: str,
    overwrite: bool | None = None,
    create_backup: bool | None = None,
) -> dict[str, Any]:
    reject_suspicious_question_marks(content, label="content")
    target = _resolve_workspace_path(path)
    _ensure_text_extension(target)

    should_overwrite = settings.default_overwrite if overwrite is None else overwrite
    should_backup = settings.create_backup if create_backup is None else create_backup

    exists = target.exists()
    if exists and not should_overwrite:
        raise FileExistsError(f"파일이 이미 존재함: {path}")
    if exists and not target.is_file():
        raise ValueError(f"파일이 아님: {path}")

    before = target.read_text(encoding="utf-8") if exists else ""
    backup_path = _backup_existing_file(target) if exists and should_backup else None

    target.parent.mkdir(parents=True, exist_ok=True)
    target.write_text(content, encoding="utf-8", newline="\n")

    return {
        "ok": True,
        "path": _to_workspace_relative(target),
        "created": not exists,
        "backup_path": backup_path,
        "diff_summary": _diff_summary(before, content),
    }


def _write_docx_from_markdown_result(
    source_path: str,
    output_path: str,
    overwrite: bool | None = None,
    create_backup: bool | None = None,
) -> dict[str, Any]:
    from docx import Document

    source = _resolve_workspace_path(source_path)
    output = _resolve_workspace_path(output_path)
    _ensure_extension(source, {".md"}, "DOCX 변환 입력 파일")
    _ensure_extension(output, {".docx"}, "DOCX 변환 출력 파일")

    if not source.exists():
        raise FileNotFoundError(f"파일을 찾을 수 없음: {source_path}")
    if not source.is_file():
        raise ValueError(f"파일이 아님: {source_path}")

    should_overwrite = settings.default_overwrite if overwrite is None else overwrite
    should_backup = settings.create_backup if create_backup is None else create_backup

    exists = output.exists()
    if exists and not should_overwrite:
        raise FileExistsError(f"파일이 이미 존재함: {output_path}")
    if exists and not output.is_file():
        raise ValueError(f"파일이 아님: {output_path}")

    markdown = source.read_text(encoding="utf-8")
    document = Document()
    style_profile = apply_docx_style(document)
    heading_count = 0
    paragraph_count = 0

    for raw_line in markdown.splitlines():
        line = raw_line.rstrip()
        stripped = line.strip()
        if not stripped:
            continue

        if stripped.startswith("#"):
            marker, _, heading = stripped.partition(" ")
            if marker and set(marker) == {"#"} and heading.strip():
                level = min(max(len(marker), 1), 4)
                document.add_heading(heading.strip(), level=level)
                heading_count += 1
                continue

        if stripped.startswith("- "):
            document.add_paragraph(stripped[2:].strip(), style="List Bullet")
            paragraph_count += 1
            continue

        document.add_paragraph(stripped)
        paragraph_count += 1

    backup_path = _backup_existing_file(output) if exists and should_backup else None
    output.parent.mkdir(parents=True, exist_ok=True)
    document.save(output)

    return {
        "ok": True,
        "path": _to_workspace_relative(output),
        "source_path": _to_workspace_relative(source),
        "output_path": _to_workspace_relative(output),
        "created": not exists,
        "backup_path": backup_path,
        "heading_count": heading_count,
        "paragraph_count": paragraph_count,
        "style_profile": style_profile,
    }


def _extract_docx_text_result(
    source_path: str,
    max_chars: int | None = None,
    include_paragraphs: bool = True,
) -> dict[str, Any]:
    from docx import Document

    source = _resolve_workspace_path(source_path)
    _ensure_extension(source, {".docx"}, "DOCX 텍스트 추출 입력 파일")

    if not source.exists():
        raise FileNotFoundError(f"파일을 찾을 수 없음: {source_path}")
    if not source.is_file():
        raise ValueError(f"파일이 아님: {source_path}")

    document = Document(source)
    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]

    table_rows: list[str] = []
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                table_rows.append(" | ".join(cells))

    blocks = paragraphs + table_rows
    full_text = "\n\n".join(blocks)
    limit = max_chars or settings.max_read_chars
    truncated = len(full_text) > limit

    return {
        "ok": True,
        "path": _to_workspace_relative(source),
        "source_path": _to_workspace_relative(source),
        "text": full_text[:limit] if truncated else full_text,
        "paragraphs": paragraphs if include_paragraphs else [],
        "paragraph_count": len(paragraphs),
        "table_row_count": len(table_rows),
        "char_count": len(full_text),
        "truncated": truncated,
    }


def _normalize_sheet_name(value: Any, fallback: str) -> str:
    raw = str(value or fallback).strip() or fallback
    for char in ("\\", "/", "*", "[", "]", ":", "?"):
        raw = raw.replace(char, " ")
    return raw[:31].strip() or fallback


def _normalize_xlsx_row(row: Any, headers: list[str]) -> list[Any]:
    if isinstance(row, dict):
        if headers:
            return [row.get(header, "") for header in headers]
        return list(row.values())
    if isinstance(row, (list, tuple)):
        return list(row)
    return [row]


def _write_xlsx_from_sheets_result(
    output_path: str,
    sheets: list[dict[str, Any]],
    overwrite: bool | None = None,
    create_backup: bool | None = None,
) -> dict[str, Any]:
    from openpyxl import Workbook
    from openpyxl.styles import Font

    reject_suspicious_question_marks(sheets, label="sheets")

    if not sheets:
        raise ValueError("sheets 값은 비어 있을 수 없음")

    output = _resolve_workspace_path(output_path)
    _ensure_extension(output, {".xlsx"}, "XLSX 출력 파일")

    should_overwrite = settings.default_overwrite if overwrite is None else overwrite
    should_backup = settings.create_backup if create_backup is None else create_backup

    exists = output.exists()
    if exists and not should_overwrite:
        raise FileExistsError(f"파일이 이미 존재함: {output_path}")
    if exists and not output.is_file():
        raise ValueError(f"파일이 아님: {output_path}")

    workbook = Workbook()
    workbook.remove(workbook.active)
    used_names: set[str] = set()
    total_rows = 0

    for index, sheet_spec in enumerate(sheets, start=1):
        base_name = _normalize_sheet_name(sheet_spec.get("name"), f"Sheet{index}")
        sheet_name = base_name
        counter = 2
        while sheet_name in used_names:
            suffix = f" {counter}"
            sheet_name = f"{base_name[:31 - len(suffix)]}{suffix}".strip()
            counter += 1
        used_names.add(sheet_name)

        worksheet = workbook.create_sheet(title=sheet_name)
        headers = [str(item) for item in sheet_spec.get("headers", [])]
        rows = sheet_spec.get("rows", [])

        if headers:
            worksheet.append(headers)
            total_rows += 1
            for cell in worksheet[1]:
                cell.font = Font(bold=True)
            worksheet.freeze_panes = "A2"

        for row in rows:
            worksheet.append(_normalize_xlsx_row(row, headers))
            total_rows += 1

        for column_cells in worksheet.columns:
            max_length = max(len(str(cell.value or "")) for cell in column_cells)
            worksheet.column_dimensions[column_cells[0].column_letter].width = min(max(max_length + 2, 10), 40)

    backup_path = _backup_existing_file(output) if exists and should_backup else None
    output.parent.mkdir(parents=True, exist_ok=True)
    workbook.save(output)

    return {
        "ok": True,
        "path": _to_workspace_relative(output),
        "output_path": _to_workspace_relative(output),
        "created": not exists,
        "backup_path": backup_path,
        "sheet_count": len(sheets),
        "row_count": total_rows,
    }


def _extract_xlsx_text_result(
    source_path: str,
    query: str | None = None,
    max_cells: int = 5000,
    max_chars: int | None = None,
) -> dict[str, Any]:
    from openpyxl import load_workbook

    source = _resolve_workspace_path(source_path)
    _ensure_extension(source, {".xlsx"}, "XLSX 텍스트 추출 입력 파일")

    if not source.exists():
        raise FileNotFoundError(f"파일을 찾을 수 없음: {source_path}")
    if not source.is_file():
        raise ValueError(f"파일이 아님: {source_path}")

    cell_limit = max(max_cells, 1)
    char_limit = max_chars or settings.max_read_chars
    normalized_query = query.casefold() if query else None

    workbook = load_workbook(source, read_only=True, data_only=True)
    sheets: list[dict[str, Any]] = []
    matches: list[dict[str, Any]] = []
    text_lines: list[str] = []
    total_cells = 0
    truncated = False

    try:
        for worksheet in workbook.worksheets:
            sheet_values: list[dict[str, Any]] = []

            for row in worksheet.iter_rows():
                for cell in row:
                    if cell.value is None:
                        continue

                    value = str(cell.value)
                    entry = {
                        "coordinate": cell.coordinate,
                        "row": cell.row,
                        "column": cell.column,
                        "value": value,
                    }
                    sheet_values.append(entry)
                    text_lines.append(f"{worksheet.title}!{cell.coordinate}: {value}")
                    total_cells += 1

                    if normalized_query and normalized_query in value.casefold():
                        matches.append({"sheet": worksheet.title, **entry})

                    if total_cells >= cell_limit:
                        truncated = True
                        break
                if truncated:
                    break

            sheets.append(
                {
                    "name": worksheet.title,
                    "cell_count": len(sheet_values),
                    "cells": sheet_values,
                }
            )

            if truncated:
                break
    finally:
        workbook.close()

    full_text = "\n".join(text_lines)
    text_truncated = len(full_text) > char_limit

    return {
        "ok": True,
        "path": _to_workspace_relative(source),
        "source_path": _to_workspace_relative(source),
        "text": full_text[:char_limit] if text_truncated else full_text,
        "sheets": sheets,
        "sheet_count": len(sheets),
        "cell_count": total_cells,
        "query": query,
        "matches": matches,
        "match_count": len(matches),
        "truncated": truncated or text_truncated,
    }


def _normalize_slide_items(value: Any) -> list[str]:
    if value is None:
        return []
    if isinstance(value, str):
        return [line.strip() for line in value.splitlines() if line.strip()]
    if isinstance(value, (list, tuple)):
        return [str(item).strip() for item in value if str(item).strip()]
    return [str(value).strip()]


def _add_textbox(slide: Any, left: Any, top: Any, width: Any, height: Any, text: str, font_size: int) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    style_paragraph(paragraph, size_pt=font_size)


def _write_pptx_from_spec_result(
    output_path: str,
    title: str,
    slides: list[dict[str, Any]],
    subtitle: str | None = None,
    overwrite: bool | None = None,
    create_backup: bool | None = None,
) -> dict[str, Any]:
    from pptx import Presentation
    from pptx.util import Inches

    reject_suspicious_question_marks(title, label="title")
    reject_suspicious_question_marks(subtitle, label="subtitle")
    reject_suspicious_question_marks(slides, label="slides")

    normalized_title = title.strip()
    if not normalized_title:
        raise ValueError("title 값은 비어 있을 수 없음")
    if not slides:
        raise ValueError("slides 값은 비어 있을 수 없음")

    output = _resolve_workspace_path(output_path)
    _ensure_extension(output, {".pptx"}, "PPTX 출력 파일")

    should_overwrite = settings.default_overwrite if overwrite is None else overwrite
    should_backup = settings.create_backup if create_backup is None else create_backup

    exists = output.exists()
    if exists and not should_overwrite:
        raise FileExistsError(f"파일이 이미 존재함: {output_path}")
    if exists and not output.is_file():
        raise ValueError(f"파일이 아님: {output_path}")

    presentation = Presentation()
    style_profile = apply_pptx_style(presentation)

    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = normalized_title
    if subtitle and title_slide.placeholders:
        try:
            title_slide.placeholders[1].text = subtitle.strip()
        except IndexError:
            pass
    style_title_slide(title_slide)

    slide_count = 1
    bullet_count = 0
    note_count = 0

    for slide_spec in slides:
        slide_title = str(slide_spec.get("title", "")).strip()
        bullets = _normalize_slide_items(slide_spec.get("bullets"))
        body = str(slide_spec.get("body", "")).strip()
        notes = str(slide_spec.get("notes", "")).strip()

        if not slide_title and not bullets and not body and not notes:
            continue

        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        if slide.shapes.title:
            slide.shapes.title.text = slide_title or "Untitled"
        style_content_slide(slide)

        content = slide.placeholders[1].text_frame
        content.clear()

        if body:
            paragraph = content.paragraphs[0]
            paragraph.text = body
            style_paragraph(paragraph, size_pt=20)

        for index, bullet in enumerate(bullets):
            paragraph = content.paragraphs[0] if index == 0 and not body else content.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            style_paragraph(paragraph, size_pt=20)
            bullet_count += 1

        if notes:
            _add_textbox(
                slide,
                Inches(0.8),
                Inches(6.5),
                Inches(11.7),
                Inches(0.7),
                f"Note: {notes}",
                10,
            )
            note_count += 1

        slide_count += 1

    backup_path = _backup_existing_file(output) if exists and should_backup else None
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output)

    return {
        "ok": True,
        "path": _to_workspace_relative(output),
        "output_path": _to_workspace_relative(output),
        "created": not exists,
        "backup_path": backup_path,
        "slide_count": slide_count,
        "bullet_count": bullet_count,
        "note_count": note_count,
        "style_profile": style_profile,
    }


def _safe_result(tool_name: str, action: Any) -> dict[str, Any]:
    try:
        ensure_base_directories()
        result = action()
    except Exception as exc:  # Tool boundaries should return readable errors to ChatGPT.
        result = {"ok": False, "error": str(exc)}
    write_operation_log(tool_name, result)
    return result


def ping_tool() -> dict[str, Any]:
    def action() -> dict[str, Any]:
        return {
            "ok": True,
            "server": "local-doc-agent",
            "workspace": str(settings.workspace_root),
            "mcp_path": "/mcp",
        }

    return _safe_result("ping", action)


def list_files_tool(
    path: str = ".",
    recursive: bool = True,
    extensions: list[str] | None = None,
) -> dict[str, Any]:
    def action() -> dict[str, Any]:
        root = _resolve_workspace_path(path)
        if not root.exists():
            return {"ok": True, "root": str(settings.workspace_root), "files": []}

        selected_extensions = extensions or DEFAULT_LIST_EXTENSIONS
        normalized_extensions = {item.lower() for item in selected_extensions}

        candidates: list[Path]
        if root.is_file():
            candidates = [root]
        elif recursive:
            candidates = [item for item in root.rglob("*") if item.is_file()]
        else:
            candidates = [item for item in root.iterdir() if item.is_file()]

        files = [
            _to_workspace_relative(item)
            for item in sorted(candidates)
            if not normalized_extensions or item.suffix.lower() in normalized_extensions
        ]
        return {"ok": True, "root": str(settings.workspace_root), "files": files}

    return _safe_result("list_files", action)


def read_text_file_tool(path: str, max_chars: int | None = None) -> dict[str, Any]:
    def action() -> dict[str, Any]:
        target = _resolve_workspace_path(path)
        _ensure_text_extension(target)
        if not target.exists():
            raise FileNotFoundError(f"파일을 찾을 수 없음: {path}")
        if not target.is_file():
            raise ValueError(f"파일이 아님: {path}")

        limit = max_chars or settings.max_read_chars
        content = target.read_text(encoding="utf-8")
        truncated = len(content) > limit
        if truncated:
            content = content[:limit]
        return {
            "ok": True,
            "path": _to_workspace_relative(target),
            "content": content,
            "truncated": truncated,
        }

    return _safe_result("read_text_file", action)


def write_text_file_tool(
    path: str,
    content: str,
    overwrite: bool | None = None,
    create_backup: bool | None = None,
) -> dict[str, Any]:
    def action() -> dict[str, Any]:
        return _write_text_file_result(
            path=path,
            content=content,
            overwrite=overwrite,
            create_backup=create_backup,
        )

    return _safe_result("write_text_file", action)


def create_markdown_tool(
    path: str,
    title: str,
    summary: str | None = None,
    sections: list[dict[str, Any]] | None = None,
    overwrite: bool | None = None,
    create_backup: bool | None = None,
) -> dict[str, Any]:
    def action() -> dict[str, Any]:
        content = _build_markdown_content(title=title, summary=summary, sections=sections)
        result = _write_text_file_result(
            path=path,
            content=content,
            overwrite=overwrite,
            create_backup=create_backup,
        )
        result["tool"] = "create_markdown"
        return result

    return _safe_result("create_markdown", action)


def export_docx_from_markdown_tool(
    source_path: str,
    output_path: str,
    overwrite: bool | None = None,
    create_backup: bool | None = None,
) -> dict[str, Any]:
    def action() -> dict[str, Any]:
        return _write_docx_from_markdown_result(
            source_path=source_path,
            output_path=output_path,
            overwrite=overwrite,
            create_backup=create_backup,
        )

    return _safe_result("export_docx_from_markdown", action)


def extract_docx_text_tool(
    source_path: str,
    max_chars: int | None = None,
    include_paragraphs: bool = True,
) -> dict[str, Any]:
    def action() -> dict[str, Any]:
        return _extract_docx_text_result(
            source_path=source_path,
            max_chars=max_chars,
            include_paragraphs=include_paragraphs,
        )

    return _safe_result("extract_docx_text", action)


def create_xlsx_from_sheets_tool(
    output_path: str,
    sheets: list[dict[str, Any]],
    overwrite: bool | None = None,
    create_backup: bool | None = None,
) -> dict[str, Any]:
    def action() -> dict[str, Any]:
        return _write_xlsx_from_sheets_result(
            output_path=output_path,
            sheets=sheets,
            overwrite=overwrite,
            create_backup=create_backup,
        )

    return _safe_result("create_xlsx_from_sheets", action)


def extract_xlsx_text_tool(
    source_path: str,
    query: str | None = None,
    max_cells: int = 5000,
    max_chars: int | None = None,
) -> dict[str, Any]:
    def action() -> dict[str, Any]:
        return _extract_xlsx_text_result(
            source_path=source_path,
            query=query,
            max_cells=max_cells,
            max_chars=max_chars,
        )

    return _safe_result("extract_xlsx_text", action)


def create_pptx_from_spec_tool(
    output_path: str,
    title: str,
    slides: list[dict[str, Any]],
    subtitle: str | None = None,
    overwrite: bool | None = None,
    create_backup: bool | None = None,
) -> dict[str, Any]:
    def action() -> dict[str, Any]:
        return _write_pptx_from_spec_result(
            output_path=output_path,
            title=title,
            slides=slides,
            subtitle=subtitle,
            overwrite=overwrite,
            create_backup=create_backup,
        )

    return _safe_result("create_pptx_from_spec", action)


def patch_text_file_tool(
    path: str,
    replacements: list[dict[str, Any]],
    create_backup: bool | None = None,
) -> dict[str, Any]:
    def action() -> dict[str, Any]:
        target = _resolve_workspace_path(path)
        _ensure_text_extension(target)
        if not target.exists():
            raise FileNotFoundError(f"파일을 찾을 수 없음: {path}")
        if not target.is_file():
            raise ValueError(f"파일이 아님: {path}")

        before = target.read_text(encoding="utf-8")
        after = before
        changed_count = 0

        for replacement in replacements:
            find = str(replacement.get("find", ""))
            replace = str(replacement.get("replace", ""))
            must_match_once = bool(replacement.get("must_match_once", False))
            reject_suspicious_question_marks(replace, label="replacements.replace")

            if find == "":
                raise ValueError("find 값은 비어 있을 수 없음")

            occurrences = after.count(find)
            if must_match_once and occurrences != 1:
                raise ValueError(f"must_match_once 실패: '{find}' 발견 횟수 {occurrences}")
            if occurrences == 0:
                continue

            after = after.replace(find, replace)
            changed_count += occurrences

        should_backup = settings.create_backup if create_backup is None else create_backup
        backup_path = _backup_existing_file(target) if changed_count > 0 and should_backup else None

        if changed_count > 0:
            target.write_text(after, encoding="utf-8", newline="\n")

        return {
            "ok": True,
            "path": _to_workspace_relative(target),
            "backup_path": backup_path,
            "changed_count": changed_count,
            "diff_summary": _diff_summary(before, after),
        }

    return _safe_result("patch_text_file", action)
