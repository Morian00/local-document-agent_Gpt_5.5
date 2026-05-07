from __future__ import annotations

import base64
import binascii
import difflib
import os
import shutil
from datetime import datetime
from pathlib import Path
from typing import Any

from .config import ensure_base_directories, settings
from .logging_utils import write_operation_log


IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp"}
IMAGE_MIME_TYPES = {
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".webp": "image/webp",
}


def _to_workspace_relative(path: Path) -> str:
    return path.resolve().relative_to(settings.workspace_root).as_posix()


def _resolve_workspace_path(path: str | None) -> Path:
    raw = "." if path is None or path == "" else path
    candidate = (settings.workspace_root / raw).resolve()
    try:
        candidate.relative_to(settings.workspace_root)
    except ValueError as exc:
        raise ValueError(f"workspace 밖의 경로는 사용할 수 없음: {path}") from exc
    return candidate


def _ensure_extension(path: Path, allowed_extensions: set[str], label: str) -> None:
    if path.suffix.lower() not in allowed_extensions:
        allowed = ", ".join(sorted(allowed_extensions))
        raise ValueError(f"{label}은 다음 확장자만 지원함: {allowed}")


def _diff_summary(before: str, after: str) -> dict[str, int]:
    before_lines = before.splitlines(keepends=True)
    after_lines = after.splitlines(keepends=True)
    diff = difflib.unified_diff(before_lines, after_lines, lineterm="")
    added = 0
    removed = 0
    for line in diff:
        if line.startswith("+++") or line.startswith("---"):
            continue
        if line.startswith("+"):
            added += 1
        elif line.startswith("-"):
            removed += 1
    return {"added_lines": added, "removed_lines": removed}


def _backup_existing_file(path: Path) -> str | None:
    if not path.exists() or not path.is_file():
        return None

    relative = path.relative_to(settings.workspace_root)
    backup_dir = settings.backups_root / relative.parent
    backup_dir.mkdir(parents=True, exist_ok=True)
    timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    backup_name = f"{path.stem}.{timestamp}{path.suffix}"
    backup_path = backup_dir / backup_name

    counter = 1
    while backup_path.exists():
        backup_path = backup_dir / f"{path.stem}.{timestamp}-{counter}{path.suffix}"
        counter += 1

    shutil.copy2(path, backup_path)
    return _to_workspace_relative(backup_path)


def _safe_result(tool_name: str, action: Any) -> dict[str, Any]:
    try:
        ensure_base_directories()
        result = action()
    except Exception as exc:  # Tool boundaries should return readable errors to ChatGPT.
        result = {"ok": False, "error": str(exc)}
    write_operation_log(tool_name, result)
    return result


def _decode_base64_image(value: str) -> bytes:
    raw = value.strip()
    if "," in raw and raw.lower().startswith("data:"):
        raw = raw.split(",", 1)[1]
    try:
        return base64.b64decode(raw, validate=True)
    except binascii.Error as exc:
        raise ValueError("image_base64 값이 올바른 base64 형식이 아님") from exc


def _png_size(data: bytes) -> tuple[int, int] | None:
    if len(data) < 24 or data[:8] != b"\x89PNG\r\n\x1a\n":
        return None
    return int.from_bytes(data[16:20], "big"), int.from_bytes(data[20:24], "big")


def _jpeg_size(data: bytes) -> tuple[int, int] | None:
    if len(data) < 4 or data[:2] != b"\xff\xd8":
        return None

    index = 2
    while index + 9 < len(data):
        if data[index] != 0xFF:
            index += 1
            continue
        marker = data[index + 1]
        index += 2
        if marker in {0xD8, 0xD9}:
            continue
        if index + 2 > len(data):
            return None
        length = int.from_bytes(data[index : index + 2], "big")
        if length < 2 or index + length > len(data):
            return None
        if 0xC0 <= marker <= 0xCF and marker not in {0xC4, 0xC8, 0xCC}:
            height = int.from_bytes(data[index + 3 : index + 5], "big")
            width = int.from_bytes(data[index + 5 : index + 7], "big")
            return width, height
        index += length
    return None


def _webp_size(data: bytes) -> tuple[int, int] | None:
    if len(data) < 30 or data[:4] != b"RIFF" or data[8:12] != b"WEBP":
        return None
    chunk = data[12:16]
    if chunk == b"VP8X" and len(data) >= 30:
        width = int.from_bytes(data[24:27], "little") + 1
        height = int.from_bytes(data[27:30], "little") + 1
        return width, height
    return None


def _image_size(data: bytes, suffix: str) -> tuple[int, int] | None:
    if suffix == ".png":
        return _png_size(data)
    if suffix in {".jpg", ".jpeg"}:
        return _jpeg_size(data)
    if suffix == ".webp":
        return _webp_size(data)
    return None


def _read_image_result(
    path: str,
    include_base64: bool = True,
    include_data_uri: bool = True,
    max_bytes: int = 2_000_000,
) -> dict[str, Any]:
    image = _resolve_workspace_path(path)
    _ensure_extension(image, IMAGE_EXTENSIONS, "이미지 읽기 파일")

    if not image.exists():
        raise FileNotFoundError(f"파일을 찾을 수 없음: {path}")
    if not image.is_file():
        raise ValueError(f"파일이 아님: {path}")

    data = image.read_bytes()
    byte_count = len(data)
    if byte_count > max_bytes:
        raise ValueError(f"이미지 파일 크기가 max_bytes를 초과함: {byte_count} > {max_bytes}")

    suffix = image.suffix.lower()
    mime_type = IMAGE_MIME_TYPES[suffix]
    encoded = base64.b64encode(data).decode("ascii") if include_base64 or include_data_uri else None
    size = _image_size(data, suffix)

    result: dict[str, Any] = {
        "ok": True,
        "path": _to_workspace_relative(image),
        "mime_type": mime_type,
        "byte_count": byte_count,
        "width": size[0] if size else None,
        "height": size[1] if size else None,
    }
    if include_base64:
        result["image_base64"] = encoded
    if include_data_uri:
        result["data_uri"] = f"data:{mime_type};base64,{encoded}"
    return result


def _write_base64_image_result(
    output_path: str,
    image_base64: str,
    overwrite: bool | None = None,
    create_backup: bool | None = None,
) -> dict[str, Any]:
    output = _resolve_workspace_path(output_path)
    _ensure_extension(output, IMAGE_EXTENSIONS, "이미지 출력 파일")

    should_overwrite = settings.default_overwrite if overwrite is None else overwrite
    should_backup = settings.create_backup if create_backup is None else create_backup

    exists = output.exists()
    if exists and not should_overwrite:
        raise FileExistsError(f"파일이 이미 존재함: {output_path}")
    if exists and not output.is_file():
        raise ValueError(f"파일이 아님: {output_path}")

    image_bytes = _decode_base64_image(image_base64)
    if not image_bytes:
        raise ValueError("image_base64 값은 비어 있을 수 없음")

    backup_path = _backup_existing_file(output) if exists and should_backup else None
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_bytes(image_bytes)

    return {
        "ok": True,
        "path": _to_workspace_relative(output),
        "output_path": _to_workspace_relative(output),
        "created": not exists,
        "backup_path": backup_path,
        "byte_count": len(image_bytes),
    }


def _insert_image_to_markdown_result(
    markdown_path: str,
    image_path: str,
    alt_text: str = "",
    position: str = "append",
    create_backup: bool | None = None,
) -> dict[str, Any]:
    target = _resolve_workspace_path(markdown_path)
    image = _resolve_workspace_path(image_path)
    _ensure_extension(target, {".md"}, "Markdown 이미지 삽입 대상")
    _ensure_extension(image, IMAGE_EXTENSIONS, "Markdown 삽입 이미지")

    if not target.exists():
        raise FileNotFoundError(f"파일을 찾을 수 없음: {markdown_path}")
    if not target.is_file():
        raise ValueError(f"파일이 아님: {markdown_path}")
    if not image.exists():
        raise FileNotFoundError(f"파일을 찾을 수 없음: {image_path}")
    if not image.is_file():
        raise ValueError(f"파일이 아님: {image_path}")

    before = target.read_text(encoding="utf-8")
    relative_image = os.path.relpath(image, target.parent).replace("\\", "/")
    image_line = f"![{alt_text.strip()}]({relative_image})"
    normalized_position = position.strip().lower()

    if normalized_position == "prepend":
        after = f"{image_line}\n\n{before.lstrip()}"
    elif normalized_position == "append":
        after = before.rstrip() + f"\n\n{image_line}\n"
    else:
        raise ValueError("position 값은 append 또는 prepend만 지원함")

    should_backup = settings.create_backup if create_backup is None else create_backup
    backup_path = _backup_existing_file(target) if should_backup else None
    target.write_text(after, encoding="utf-8", newline="\n")

    return {
        "ok": True,
        "path": _to_workspace_relative(target),
        "image_path": _to_workspace_relative(image),
        "backup_path": backup_path,
        "position": normalized_position,
        "diff_summary": _diff_summary(before, after),
    }


def _insert_image_to_pptx_result(
    pptx_path: str,
    image_path: str,
    slide_index: int = 0,
    left: float = 1.0,
    top: float = 1.5,
    width: float | None = 4.0,
    height: float | None = None,
    create_backup: bool | None = None,
) -> dict[str, Any]:
    from pptx import Presentation
    from pptx.util import Inches

    target = _resolve_workspace_path(pptx_path)
    image = _resolve_workspace_path(image_path)
    _ensure_extension(target, {".pptx"}, "PPTX 이미지 삽입 대상")
    _ensure_extension(image, IMAGE_EXTENSIONS, "PPTX 삽입 이미지")

    if not target.exists():
        raise FileNotFoundError(f"파일을 찾을 수 없음: {pptx_path}")
    if not target.is_file():
        raise ValueError(f"파일이 아님: {pptx_path}")
    if not image.exists():
        raise FileNotFoundError(f"파일을 찾을 수 없음: {image_path}")
    if not image.is_file():
        raise ValueError(f"파일이 아님: {image_path}")

    presentation = Presentation(target)
    if slide_index < 0 or slide_index >= len(presentation.slides):
        raise ValueError(f"slide_index 범위 초과: {slide_index}")

    slide = presentation.slides[slide_index]
    width_value = Inches(width) if width is not None else None
    height_value = Inches(height) if height is not None else None
    slide.shapes.add_picture(
        str(image),
        Inches(left),
        Inches(top),
        width=width_value,
        height=height_value,
    )

    should_backup = settings.create_backup if create_backup is None else create_backup
    backup_path = _backup_existing_file(target) if should_backup else None
    presentation.save(target)

    return {
        "ok": True,
        "path": _to_workspace_relative(target),
        "image_path": _to_workspace_relative(image),
        "backup_path": backup_path,
        "slide_index": slide_index,
        "shape_count": len(slide.shapes),
    }


def list_assets_tool(
    path: str = "assets",
    recursive: bool = True,
    extensions: list[str] | None = None,
) -> dict[str, Any]:
    def action() -> dict[str, Any]:
        root = _resolve_workspace_path(path)
        if not root.exists():
            return {"ok": True, "root": str(settings.workspace_root / "assets"), "files": []}

        selected_extensions = extensions or sorted(IMAGE_EXTENSIONS)
        normalized_extensions = {item.lower() for item in selected_extensions}

        candidates: list[Path]
        if root.is_file():
            candidates = [root]
        elif recursive:
            candidates = [item for item in root.rglob("*") if item.is_file()]
        else:
            candidates = [item for item in root.iterdir() if item.is_file()]

        files = [
            _to_workspace_relative(item)
            for item in sorted(candidates)
            if not normalized_extensions or item.suffix.lower() in normalized_extensions
        ]
        return {"ok": True, "root": str(root), "files": files}

    return _safe_result("list_assets", action)


def save_base64_image_tool(
    output_path: str,
    image_base64: str,
    overwrite: bool | None = None,
    create_backup: bool | None = None,
) -> dict[str, Any]:
    def action() -> dict[str, Any]:
        return _write_base64_image_result(
            output_path=output_path,
            image_base64=image_base64,
            overwrite=overwrite,
            create_backup=create_backup,
        )

    return _safe_result("save_base64_image", action)


def read_image_file_tool(
    path: str,
    include_base64: bool = True,
    include_data_uri: bool = True,
    max_bytes: int = 2_000_000,
) -> dict[str, Any]:
    def action() -> dict[str, Any]:
        return _read_image_result(
            path=path,
            include_base64=include_base64,
            include_data_uri=include_data_uri,
            max_bytes=max_bytes,
        )

    return _safe_result("read_image_file", action)


def insert_image_to_markdown_tool(
    markdown_path: str,
    image_path: str,
    alt_text: str = "",
    position: str = "append",
    create_backup: bool | None = None,
) -> dict[str, Any]:
    def action() -> dict[str, Any]:
        return _insert_image_to_markdown_result(
            markdown_path=markdown_path,
            image_path=image_path,
            alt_text=alt_text,
            position=position,
            create_backup=create_backup,
        )

    return _safe_result("insert_image_to_markdown", action)


def insert_image_to_pptx_tool(
    pptx_path: str,
    image_path: str,
    slide_index: int = 0,
    left: float = 1.0,
    top: float = 1.5,
    width: float | None = 4.0,
    height: float | None = None,
    create_backup: bool | None = None,
) -> dict[str, Any]:
    def action() -> dict[str, Any]:
        return _insert_image_to_pptx_result(
            pptx_path=pptx_path,
            image_path=image_path,
            slide_index=slide_index,
            left=left,
            top=top,
            width=width,
            height=height,
            create_backup=create_backup,
        )

    return _safe_result("insert_image_to_pptx", action)
